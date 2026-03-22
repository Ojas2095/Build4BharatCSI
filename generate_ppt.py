"""
PMDDKY Smart Governance Platform — BUILD4BHARAT Hackathon 9.0
Team: Kishta Coders
PPT Generator — Minimal text, max visual, wireframes embedded
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── Image paths (AI-generated wireframes/diagrams) ──────────────────────────
BRAIN_DIR = r"C:\Users\Dev\.gemini\antigravity\brain\67e9ac84-4b3a-48e0-8a51-5b9418f34317"

def img(name):
    """Find the latest version of a generated image by prefix."""
    matches = [f for f in os.listdir(BRAIN_DIR) if f.startswith(name) and f.endswith(".png")]
    if not matches:
        raise FileNotFoundError(f"No image found for prefix '{name}' in {BRAIN_DIR}")
    matches.sort()  # latest alphabetically (timestamp suffix)
    return os.path.join(BRAIN_DIR, matches[-1])

# ── Colours ──────────────────────────────────────────────────────────────────
DARK_GREEN   = RGBColor(0x0D, 0x3D, 0x22)
MID_GREEN    = RGBColor(0x1A, 0x5C, 0x38)
BRIGHT_GREEN = RGBColor(0x1A, 0x7C, 0x45)
ACCENT_GREEN = RGBColor(0x00, 0xC9, 0x7A)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE    = RGBColor(0xF5, 0xFA, 0xF7)
DARK_TEXT    = RGBColor(0x11, 0x11, 0x11)
LIGHT_GREY   = RGBColor(0xCC, 0xCC, 0xCC)
GOLD         = RGBColor(0xD4, 0xA8, 0x20)
BLUE         = RGBColor(0x1E, 0x90, 0xFF)
RED          = RGBColor(0xCC, 0x22, 0x22)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


# ── Helpers ──────────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill_rgb):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.line.fill.background()
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    return shape


def add_text(slide, text, x, y, w, h,
             font_size=12, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(font_size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def add_image(slide, path, x, y, w, h):
    return slide.shapes.add_picture(path, x, y, w, h)


def header_bar(slide, show_brand=True):
    """Standard dark-green header for content slides."""
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.72), DARK_GREEN)
    if show_brand:
        add_text(slide, "BUILD4भारत", Inches(0.15), Inches(0.1),
                 Inches(2.4), Inches(0.55), font_size=17, bold=True, color=WHITE)
    add_text(slide, "D❖LL Technologies", Inches(4.9), Inches(0.14),
             Inches(3.5), Inches(0.5), font_size=15, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)
    sponsors = "MeitY  |  UNDP  |  CSI  |  Anritsu  |  runway"
    add_text(slide, sponsors, Inches(9.1), Inches(0.22),
             Inches(4.1), Inches(0.35), font_size=8, color=LIGHT_GREY,
             align=PP_ALIGN.RIGHT)


def slide_footer(slide, num, total=9):
    add_text(slide, f"SLIDE {num} OF {total}",
             Inches(11.5), Inches(7.25), Inches(1.8), Inches(0.25),
             font_size=8, color=RGBColor(0x88, 0x88, 0x88),
             align=PP_ALIGN.RIGHT)


def slide_tag(slide, text):
    """Big section label bottom-left of header."""
    add_text(slide, text, Inches(0.15), Inches(0.72),
             Inches(13), Inches(0.38), font_size=11, bold=True,
             color=ACCENT_GREEN, align=PP_ALIGN.LEFT)


def add_stat_badge(slide, label, value, x, y, w=Inches(2.8), h=Inches(1.1)):
    """A green-bordered stat card."""
    add_rect(slide, x, y, w, h, WHITE)
    add_rect(slide, x, y, w, Inches(0.06), BRIGHT_GREEN)
    add_text(slide, value, x + Inches(0.1), y + Inches(0.08),
             w - Inches(0.15), Inches(0.55), font_size=24, bold=True,
             color=DARK_GREEN, align=PP_ALIGN.LEFT)
    add_text(slide, label, x + Inches(0.1), y + Inches(0.62),
             w - Inches(0.15), Inches(0.38), font_size=9, color=DARK_TEXT,
             align=PP_ALIGN.LEFT)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK)

# Full dark background
add_rect(s1, 0, 0, SLIDE_W, SLIDE_H, DARK_GREEN)
add_rect(s1, 0, 0, SLIDE_W, SLIDE_H, MID_GREEN)

# Top header — title slide variant
add_rect(s1, 0, 0, SLIDE_W, Inches(0.7), DARK_GREEN)
add_text(s1, "SCHOOL OF COMPUTER SCIENCE", Inches(0.15), Inches(0.1),
         Inches(2.2), Inches(0.55), font_size=8, color=LIGHT_GREY)
add_text(s1, "MeitY Startup Hub  |  UNDP  |  CSI  |  Anritsu  |  runway",
         Inches(2.8), Inches(0.18), Inches(7.2), Inches(0.4),
         font_size=8, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s1, "©UPES", Inches(12.0), Inches(0.12),
         Inches(1.2), Inches(0.5), font_size=18, bold=True, color=WHITE,
         align=PP_ALIGN.RIGHT)

# Hackathon branding
add_text(s1, "BUILD4भारत", Inches(1.0), Inches(0.95),
         Inches(5.5), Inches(0.8), font_size=38, bold=True, color=WHITE)
add_text(s1, "HACKATHON 9.0", Inches(1.0), Inches(1.75),
         Inches(5.5), Inches(0.55), font_size=24, bold=True, color=WHITE)
add_text(s1, "Hack for Humanity  ·  Build for Community",
         Inches(1.0), Inches(2.3), Inches(5.5), Inches(0.4),
         font_size=9, color=RGBColor(0xAA, 0xD4, 0xBC))
add_text(s1, "Technical Partner:  D❖LL Technologies",
         Inches(1.0), Inches(2.7), Inches(5.5), Inches(0.38),
         font_size=11, bold=True, color=WHITE)

# Vertical divider
add_rect(s1, Inches(7.0), Inches(0.9), Inches(0.03), Inches(5.8), BRIGHT_GREEN)

# Team white card (right side)
card_x, card_y, card_w, card_h = Inches(7.4), Inches(1.1), Inches(5.7), Inches(5.4)
add_rect(s1, card_x, card_y, card_w, card_h, WHITE)
add_rect(s1, card_x, card_y, Inches(0.06), card_h, BRIGHT_GREEN)

add_text(s1, "KISHTA CODERS", card_x + Inches(0.2), card_y + Inches(0.2),
         card_w - Inches(0.3), Inches(0.55), font_size=22, bold=True, color=DARK_GREEN)
add_text(s1, "Team Name:  Kishta Coders",
         card_x + Inches(0.2), card_y + Inches(0.85),
         card_w - Inches(0.3), Inches(0.38), font_size=11, color=DARK_TEXT)
add_text(s1, "Team Leader:  ________________________________",
         card_x + Inches(0.2), card_y + Inches(1.25),
         card_w - Inches(0.3), Inches(0.38), font_size=11, color=DARK_TEXT)
add_rect(s1, card_x + Inches(0.2), card_y + Inches(1.72),
         card_w - Inches(0.4), Inches(0.02), BRIGHT_GREEN)
add_text(s1, "PROBLEM STATEMENT",
         card_x + Inches(0.2), card_y + Inches(1.85),
         card_w - Inches(0.3), Inches(0.32), font_size=9, bold=True,
         color=BRIGHT_GREEN)
add_text(s1,
         "Real-time Smart Governance for PMDDKY\n"
         "Live Monitoring & AI-powered Analytics\n"
         "Dashboard for Uttarakhand's 13 Districts",
         card_x + Inches(0.2), card_y + Inches(2.2),
         card_w - Inches(0.3), Inches(1.0),
         font_size=13, bold=True, color=DARK_GREEN, wrap=True)

# Three mini stat badges on the card bottom
stats = [("5.8L+", "Farmers"), ("13", "Districts"), ("₹847Cr", "Target Funds")]
for i, (val, lbl) in enumerate(stats):
    bx = card_x + Inches(0.2) + i * Inches(1.78)
    by = card_y + Inches(3.55)
    add_rect(s1, bx, by, Inches(1.65), Inches(0.85), OFF_WHITE)
    add_text(s1, val, bx + Inches(0.1), by + Inches(0.05),
             Inches(1.4), Inches(0.45), font_size=18, bold=True, color=DARK_GREEN)
    add_text(s1, lbl, bx + Inches(0.1), by + Inches(0.5),
             Inches(1.4), Inches(0.28), font_size=8, color=DARK_TEXT)

slide_footer(s1, 1)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — PROBLEM STATEMENT  (full-bleed infographic)
# ════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK)
add_rect(s2, 0, 0, SLIDE_W, SLIDE_H, RGBColor(0xF8, 0xFB, 0xF9))
header_bar(s2)

# Slide label
add_text(s2, "PROBLEM STATEMENT", Inches(0.3), Inches(0.78),
         Inches(6), Inches(0.38), font_size=14, bold=True, color=DARK_GREEN)

# Big hero stat
add_rect(s2, Inches(0.2), Inches(1.25), Inches(3.4), Inches(1.55), DARK_GREEN)
add_text(s2, "5.8 LAKH", Inches(0.3), Inches(1.3),
         Inches(3.2), Inches(0.75), font_size=36, bold=True, color=ACCENT_GREEN,
         align=PP_ALIGN.CENTER)
add_text(s2, "Farmers eligible in\nUttarakhand under PMDDKY",
         Inches(0.3), Inches(2.0), Inches(3.2), Inches(0.65),
         font_size=10, color=WHITE, align=PP_ALIGN.CENTER)

add_rect(s2, Inches(0.2), Inches(2.95), Inches(3.4), Inches(1.0), MID_GREEN)
add_text(s2, "13 Backward Districts", Inches(0.3), Inches(3.0),
         Inches(3.2), Inches(0.38), font_size=14, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER)
add_text(s2, "Zero unified visibility", Inches(0.3), Inches(3.38),
         Inches(3.2), Inches(0.38), font_size=9,
         color=RGBColor(0xAA, 0xDD, 0xBB), align=PP_ALIGN.CENTER)

add_rect(s2, Inches(0.2), Inches(4.1), Inches(3.4), Inches(0.9), RED)
add_text(s2, "MONTHS of delay", Inches(0.3), Inches(4.15),
         Inches(3.2), Inches(0.38), font_size=14, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER)
add_text(s2, "in fund disbursement", Inches(0.3), Inches(4.52),
         Inches(3.2), Inches(0.3), font_size=9,
         color=RGBColor(0xFF, 0xCC, 0xCC), align=PP_ALIGN.CENTER)

# Full infographic image (problem viz)
try:
    add_image(s2, img("problem_infographic"),
              Inches(3.8), Inches(1.1), Inches(9.3), Inches(4.2))
except Exception as e:
    add_text(s2, f"[Image: problem_infographic — {e}]",
             Inches(3.8), Inches(1.1), Inches(9.3), Inches(4.2),
             font_size=10, color=DARK_TEXT)

# Caption
add_text(s2,
         "District officials rely on paper registers & siloed Excel files — no central platform, "
         "no real-time visibility, no anomaly detection",
         Inches(0.2), Inches(5.2), Inches(12.9), Inches(0.4),
         font_size=9, italic=True, color=RGBColor(0x55, 0x55, 0x55))

slide_footer(s2, 2)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — PROPOSED SOLUTION  (dashboard wireframe)
# ════════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(BLANK)
add_rect(s3, 0, 0, SLIDE_W, SLIDE_H, RGBColor(0xF4, 0xF9, 0xF6))
header_bar(s3)

add_text(s3, "PROPOSED SOLUTION", Inches(0.3), Inches(0.78),
         Inches(8), Inches(0.38), font_size=14, bold=True, color=DARK_GREEN)

# Left: 3 solution pillars
pillars = [
    ("📡", "Real-Time Dashboard", "Live KPI cards — beneficiaries, funds, district progress, % verified"),
    ("🗺️", "GIS District Map",   "Choropleth map of Uttarakhand's 13 districts — click to drill down"),
    ("🤖", "AI Analytics",       "Anomaly detection, trend forecasting, district performance ranking"),
]
py = Inches(1.25)
for icon, title, desc in pillars:
    add_rect(s3, Inches(0.18), py, Inches(4.0), Inches(1.2), WHITE)
    add_rect(s3, Inches(0.18), py, Inches(0.06), Inches(1.2), BRIGHT_GREEN)
    add_text(s3, icon + "  " + title, Inches(0.32), py + Inches(0.08),
             Inches(3.7), Inches(0.4), font_size=11, bold=True, color=DARK_GREEN)
    add_text(s3, desc, Inches(0.32), py + Inches(0.5),
             Inches(3.7), Inches(0.6), font_size=9, color=DARK_TEXT, wrap=True)
    py += Inches(1.32)

# Two more pillars below
pillars2 = [
    ("📋", "Beneficiary\nManagement", "Aadhaar-verified records, status badges"),
    ("⚡", "Socket.IO\nReal-Time", "Live push — no page reload needed"),
]
for j, (icon, title, desc) in enumerate(pillars2):
    bx = Inches(0.18) + j * Inches(2.1)
    by = py
    add_rect(s3, bx, by, Inches(2.0), Inches(1.1), WHITE)
    add_rect(s3, bx, by, Inches(0.06), Inches(1.1), ACCENT_GREEN)
    add_text(s3, icon + " " + title, bx + Inches(0.14), by + Inches(0.08),
             Inches(1.78), Inches(0.55), font_size=10, bold=True, color=DARK_GREEN)
    add_text(s3, desc, bx + Inches(0.14), by + Inches(0.65),
             Inches(1.78), Inches(0.38), font_size=8, color=DARK_TEXT, wrap=True)

# Right: dashboard wireframe (big)
try:
    add_image(s3, img("dashboard_wireframe"),
              Inches(4.4), Inches(0.85), Inches(8.75), Inches(5.55))
except Exception as e:
    add_text(s3, f"[Dashboard wireframe — {e}]",
             Inches(4.4), Inches(0.85), Inches(8.75), Inches(5.55),
             font_size=10, color=DARK_TEXT)

add_rect(s3, Inches(4.4), Inches(6.4), Inches(8.75), Inches(0.26), DARK_GREEN)
add_text(s3, "▶  Dashboard Wireframe — PMDDKY Smart Governance Platform (React 18 + Leaflet + Socket.IO)",
         Inches(4.55), Inches(6.43), Inches(8.5), Inches(0.22),
         font_size=8, color=ACCENT_GREEN)

slide_footer(s3, 3)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — TECHNOLOGY STACK
# ════════════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(BLANK)
add_rect(s4, 0, 0, SLIDE_W, SLIDE_H, RGBColor(0xF8, 0xFB, 0xF9))
header_bar(s4)

add_text(s4, "TECHNOLOGY STACK", Inches(0.3), Inches(0.78),
         Inches(8), Inches(0.38), font_size=14, bold=True, color=DARK_GREEN)

# Left narrow: why this stack
why = [
    "100% Open Source\nZero licensing cost",
    "Deployable on ₹500/mo VPS\nor govt NIC servers",
    "Handles concurrent\ndistrict-level users",
    "Real-time via\nSocket.IO",
]
wy = Inches(1.2)
for w in why:
    add_rect(s4, Inches(0.15), wy, Inches(2.55), Inches(0.85), WHITE)
    add_rect(s4, Inches(0.15), wy, Inches(0.05), Inches(0.85), BRIGHT_GREEN)
    add_text(s4, w, Inches(0.28), wy + Inches(0.08),
             Inches(2.3), Inches(0.72), font_size=9, color=DARK_TEXT, wrap=True)
    wy += Inches(0.98)

# Right: tech stack visual diagram
try:
    add_image(s4, img("tech_stack_diagram"),
              Inches(2.9), Inches(0.85), Inches(10.25), Inches(5.6))
except Exception as e:
    add_text(s4, f"[Tech stack diagram — {e}]",
             Inches(2.9), Inches(0.85), Inches(10.25), Inches(5.6),
             font_size=10, color=DARK_TEXT)

slide_footer(s4, 4)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — METHODOLOGY / SYSTEM ARCHITECTURE
# ════════════════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(BLANK)
add_rect(s5, 0, 0, SLIDE_W, SLIDE_H, RGBColor(0xF4, 0xF9, 0xF6))
header_bar(s5)

add_text(s5, "METHODOLOGY & SYSTEM ARCHITECTURE", Inches(0.3), Inches(0.78),
         Inches(10), Inches(0.38), font_size=14, bold=True, color=DARK_GREEN)

# Architecture diagram — big
try:
    add_image(s5, img("system_architecture"),
              Inches(0.15), Inches(1.12), Inches(8.0), Inches(4.5))
except Exception as e:
    add_text(s5, f"[System architecture — {e}]",
             Inches(0.15), Inches(1.12), Inches(8.0), Inches(4.5),
             font_size=10, color=DARK_TEXT)

# Right: data flow diagram
try:
    add_image(s5, img("data_flow_diagram"),
              Inches(8.3), Inches(1.12), Inches(4.9), Inches(2.1))
except Exception as e:
    add_text(s5, f"[Data flow — {e}]",
             Inches(8.3), Inches(1.12), Inches(4.9), Inches(2.1),
             font_size=10, color=DARK_TEXT)

# Beneficiary table wireframe (bottom-right)
try:
    add_image(s5, img("beneficiary_table_wireframe"),
              Inches(8.3), Inches(3.35), Inches(4.9), Inches(2.3))
except Exception as e:
    add_text(s5, f"[Beneficiary table — {e}]",
             Inches(8.3), Inches(3.35), Inches(4.9), Inches(2.3),
             font_size=10, color=DARK_TEXT)

# Caption strip
add_rect(s5, Inches(0.15), Inches(5.72), Inches(13.0), Inches(0.3), DARK_GREEN)
add_text(s5, "3-tier architecture: Field Input → Node.js + Socket.IO Backend → React Frontend  |  AI Engine: Anomaly Detection + Forecasting",
         Inches(0.3), Inches(5.75), Inches(12.7), Inches(0.25),
         font_size=8, color=ACCENT_GREEN)

slide_footer(s5, 5)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — IMPACT
# ════════════════════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(BLANK)
add_rect(s6, 0, 0, SLIDE_W, SLIDE_H, RGBColor(0xF8, 0xFB, 0xF9))
header_bar(s6)

add_text(s6, "IMPACT — TECH FOR COMMUNITY", Inches(0.3), Inches(0.78),
         Inches(8), Inches(0.38), font_size=14, bold=True, color=DARK_GREEN)

# Large impact infographic
try:
    add_image(s6, img("impact_infographic"),
              Inches(0.15), Inches(1.12), Inches(8.3), Inches(5.3))
except Exception as e:
    add_text(s6, f"[Impact infographic — {e}]",
             Inches(0.15), Inches(1.12), Inches(8.3), Inches(5.3),
             font_size=10, color=DARK_TEXT)

# Right column: who benefits
beneficiaries = [
    ("👨‍🌾", "5.8L+ Farmers",
     "Faster verified fund access — months → days"),
    ("🏛️", "13 District Admins",
     "Single dashboard replaces manual reports"),
    ("📈", "State Policymakers",
     "Compare all districts, reallocate proactively"),
    ("🌱", "Agriculture",
     "Better seeds & irrigation from timely fund flow"),
]
bry = Inches(1.2)
for icon, title, desc in beneficiaries:
    add_rect(s6, Inches(8.6), bry, Inches(4.55), Inches(1.05), WHITE)
    add_rect(s6, Inches(8.6), bry, Inches(0.05), Inches(1.05), BRIGHT_GREEN)
    add_text(s6, icon + "  " + title, Inches(8.73), bry + Inches(0.07),
             Inches(4.3), Inches(0.38), font_size=11, bold=True, color=DARK_GREEN)
    add_text(s6, desc, Inches(8.73), bry + Inches(0.48),
             Inches(4.3), Inches(0.45), font_size=9, color=DARK_TEXT, wrap=True)
    bry += Inches(1.18)

slide_footer(s6, 6)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — FEASIBILITY & INNOVATION
# ════════════════════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(BLANK)
add_rect(s7, 0, 0, SLIDE_W, SLIDE_H, RGBColor(0xF4, 0xF9, 0xF6))
header_bar(s7)

add_text(s7, "FEASIBILITY & INNOVATION", Inches(0.3), Inches(0.78),
         Inches(8), Inches(0.38), font_size=14, bold=True, color=DARK_GREEN)

# Comparison infographic (big)
try:
    add_image(s7, img("feasibility_comparison"),
              Inches(0.15), Inches(1.1), Inches(7.7), Inches(5.4))
except Exception as e:
    add_text(s7, f"[Feasibility comparison — {e}]",
             Inches(0.15), Inches(1.1), Inches(7.7), Inches(5.4),
             font_size=10, color=DARK_TEXT)

# Right: innovation highlights
innovations = [
    ("🗺️", "First GIS-Mapped\nPMDDKY Dashboard",
     "Choropleth map of Uttarakhand — district-level beneficiary density live"),
    ("🤖", "AI Anomaly\nDetection",
     "Flags unusual fund spikes or enrollment drops before they become crises"),
    ("⚡", "Zero-Reload\nReal-Time",
     "Socket.IO — field officers and HQ see the same data change simultaneously"),
    ("💰", "₹500/month\nDeployment",
     "100% open-source, runs on cheap VPS or government NIC servers"),
]
iy = Inches(1.15)
for icon, title, desc in innovations:
    add_rect(s7, Inches(8.05), iy, Inches(5.1), Inches(1.22), WHITE)
    add_rect(s7, Inches(8.05), iy, Inches(0.05), Inches(1.22), ACCENT_GREEN)
    add_text(s7, icon + "  " + title, Inches(8.18), iy + Inches(0.07),
             Inches(4.8), Inches(0.48), font_size=10, bold=True, color=DARK_GREEN)
    add_text(s7, desc, Inches(8.18), iy + Inches(0.58),
             Inches(4.8), Inches(0.52), font_size=9, color=DARK_TEXT, wrap=True)
    iy += Inches(1.34)

slide_footer(s7, 7)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — FUTURE SCOPE
# ════════════════════════════════════════════════════════════════════════════
s8 = prs.slides.add_slide(BLANK)
add_rect(s8, 0, 0, SLIDE_W, SLIDE_H, RGBColor(0xF8, 0xFB, 0xF9))
header_bar(s8)

add_text(s8, "FUTURE SCOPE & ROADMAP", Inches(0.3), Inches(0.78),
         Inches(8), Inches(0.38), font_size=14, bold=True, color=DARK_GREEN)

# Roadmap image (big, left)
try:
    add_image(s8, img("future_roadmap"),
              Inches(0.15), Inches(1.1), Inches(8.8), Inches(5.45))
except Exception as e:
    add_text(s8, f"[Future roadmap — {e}]",
             Inches(0.15), Inches(1.1), Inches(8.8), Inches(5.45),
             font_size=10, color=DARK_TEXT)

# Right: scope summary cards
scope = [
    ("📱", "Mobile App", "React Native offline-first for field officers"),
    ("🔗", "Aadhaar API", "Real integration for instant beneficiary verification"),
    ("🔒", "Blockchain", "Immutable audit trail for RTI compliance"),
    ("🌐", "100 Districts", "National rollout to all PMDDKY districts"),
]
sy = Inches(1.2)
for icon, title, desc in scope:
    add_rect(s8, Inches(9.18), sy, Inches(3.97), Inches(1.2), WHITE)
    add_rect(s8, Inches(9.18), sy, Inches(0.05), Inches(1.2), BRIGHT_GREEN)
    add_text(s8, icon + "  " + title, Inches(9.31), sy + Inches(0.08),
             Inches(3.7), Inches(0.38), font_size=11, bold=True, color=DARK_GREEN)
    add_text(s8, desc, Inches(9.31), sy + Inches(0.5),
             Inches(3.7), Inches(0.55), font_size=9, color=DARK_TEXT, wrap=True)
    sy += Inches(1.35)

slide_footer(s8, 8)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — THANK YOU
# ════════════════════════════════════════════════════════════════════════════
s9 = prs.slides.add_slide(BLANK)
add_rect(s9, 0, 0, SLIDE_W, SLIDE_H, DARK_GREEN)

# Green gradient overlay — simulate with two rects
add_rect(s9, 0, 0, SLIDE_W, SLIDE_H, MID_GREEN)

# Large "Thank You"
txb = s9.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(12.33), Inches(2.0))
tf9 = txb.text_frame
p = tf9.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "Thank You"
run.font.size  = Pt(72)
run.font.bold  = True
run.font.color.rgb = WHITE

add_text(s9, "KISHTA CODERS  |  BUILD4BHARAT HACKATHON 9.0",
         Inches(0.5), Inches(4.0), Inches(12.33), Inches(0.5),
         font_size=16, bold=True, color=ACCENT_GREEN, align=PP_ALIGN.CENTER)

add_text(s9, "PMDDKY Smart Governance Platform — Hack for Humanity, Build for Community",
         Inches(0.5), Inches(4.6), Inches(12.33), Inches(0.4),
         font_size=11, color=LIGHT_GREY, align=PP_ALIGN.CENTER)

# Accent line
add_rect(s9, Inches(4.0), Inches(5.2), Inches(5.33), Inches(0.04), BRIGHT_GREEN)

add_text(s9,
         "Real-time Dashboard  ·  GIS Mapping  ·  AI Anomaly Detection  ·  Socket.IO  ·  100% Open Source",
         Inches(0.5), Inches(5.4), Inches(12.33), Inches(0.35),
         font_size=9, color=RGBColor(0x88, 0x88, 0x88), align=PP_ALIGN.CENTER)

slide_footer(s9, 9)


# ── Save ────────────────────────────────────────────────────────────────────
output_path = r"E:\CSI_B4B\KishtaCoders_PMDDKY_BUILD4BHARAT.pptx"
prs.save(output_path)
print(f"✅  Saved: {output_path}")
